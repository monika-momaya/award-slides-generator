#!/usr/bin/env python3
"""
Award Slide Generator (template-agnostic)
===========================================
Reads a master Excel sheet of award categories/nominees/winners and a
PowerPoint certificate template, and auto-generates ONE combined deck: for
every award category (per zone, if that category has zone-level results),
a Nominee slide (every nominee, shuffled order) immediately followed by
that award's Winner slide (the winning company's name only) -- or just
Winner slides on their own, if that's all the template provides (see
below: every placeholder is optional).

WORKS WITH ANY YEAR'S TEMPLATE, because a slide's ROLE is decided purely by
which literal placeholder token(s) it contains -- not its position in the
file, and EVERY token is optional. Place these EXACT tokens (literal
match, case-sensitive as written) inside text boxes on your template:

  - "<<NOMINEES>>"      -> marks this slide as the Nominee stencil. The box
                           itself receives the full nominee list for that
                           award (one name per line, shuffled), keeping
                           whatever font/size/position that box already has
                           on your slide.
  - "<<WINNER>>"         -> marks this slide as the Winner stencil. The box
                           receives that award's winning company's name,
                           keeping the box's existing formatting/position.
  - "<<AWARD CATEGORY>>" -> (optional, on either stencil) receives the
                           award category title.
  - "<<ZONE>>"           -> (optional, on either stencil) receives the zone
                           name (North/South/etc), when the category is
                           zone-split. If a given award has no zone, this
                           box is removed entirely rather than left showing
                           the literal token text.
  - "<<nominees-word>>"  -> (optional, anywhere) replaced with the literal
                           word "NOMINEES" -- useful for templates that
                           want that word printed as a static heading.
  - "<<winner-word>>"    -> (optional, anywhere) replaced with the literal
                           word "WINNER", same idea.

A slide needs AT LEAST ONE of <<NOMINEES>> / <<WINNER>> to be used as a
per-award stencil at all; every other token on it is optional and simply
left as the box's own design if omitted. A template can therefore be:
  - Nominee slide + Winner slide (two stencils, the classic case), or
  - Winner slide only (one stencil) -- e.g. events with no nominee reveal,
    just a single winner announcement design, or
  - Nominee slide only.
Any other slide in the template (no <<NOMINEES>>/<<WINNER>> token at all)
is left completely untouched and copied through once at the end of the
deck, e.g. a title or thank-you slide.

Placeholder boxes are NEVER auto-repositioned or auto-resized -- each one
keeps the exact position, size, and text formatting you already gave it on
the template slide; only its text content is replaced. Design the slide
exactly the way you want it to look, and that's how it'll look.

WORKS WITH ANY YEAR'S EXCEL, regardless of column names, because columns
are identified by the SHAPE of their data rather than their header text:
  - Category column: filled only on the first row of each category block
    (sparse), longer descriptive text.
  - Company/Nominee Name column: filled on every row, mostly unique values.
  - Zone/Region column (optional): filled on every row, small repeating
    set of values (e.g. North/South/East/West).
  - Result column: filled on every row, small repeating set of values
    (e.g. Winner/1st Runnerup/2nd Runnerup).
  Within each category/zone block, the winner is identified by checking
  for the word "winner" in the result column's text; if that doesn't
  cleanly identify exactly one row, the FIRST row of the block is treated
  as the winner instead (matching the standard convention of listing the
  best result first). Header text is used only as a minor tie-breaker, so
  renaming columns next year will not break this -- but wildly reordering
  or removing data (e.g. no column with a small repeating result
  vocabulary at all) will, since at that point there's no signal left to
  detect a winner from.
A blank row separates one zone's group of nominees from the next zone's
group within the same category (matching the standard BOCI/Prawaas
export format).

Usage:
    python3 generate_award_slides.py <excel_file> <template_pptx> <output_file.pptx>
"""
import sys
import os
import re
import copy
import random
import openpyxl
from pptx import Presentation

RANDOM_SEED = 42


# --------------------------------------------------------------------------
# Excel parsing
# --------------------------------------------------------------------------

def profile_column(values_by_row, total_rows):
    """Compute structural statistics for one column's non-empty values,
    used to guess its semantic role without relying on its header text."""
    non_empty = [v for v in values_by_row if v is not None and str(v).strip() != ""]
    n = len(non_empty)
    if n == 0:
        return {"fill_ratio": 0, "uniq_ratio": 0, "avg_len": 0, "distinct": 0,
                "pct_numeric": 0, "pct_at_sign": 0, "pct_alpha_rich": 0, "n": 0}
    texts = [str(v).strip() for v in non_empty]
    distinct = len(set(t.lower() for t in texts))
    avg_len = sum(len(t) for t in texts) / n
    pct_numeric = sum(1 for t in texts if re.fullmatch(r"[\d\-+() ]+", t)) / n
    pct_at_sign = sum(1 for t in texts if "@" in t) / n
    pct_alpha_rich = sum(1 for t in texts if re.search(r"[A-Za-z]{3,}", t)) / n
    return {
        "fill_ratio": n / total_rows if total_rows else 0,
        "uniq_ratio": distinct / n,
        "avg_len": avg_len,
        "distinct": distinct,
        "pct_numeric": pct_numeric,
        "pct_at_sign": pct_at_sign,
        "pct_alpha_rich": pct_alpha_rich,
        "n": n,
    }


def header_bonus(header_text, keyword_groups):
    """Small score bonus if the header text matches one of the candidate
    keyword-sets. Purely a tie-breaker -- never required for a match."""
    if not header_text:
        return 0
    text = str(header_text).lower()
    for i, keywords in enumerate(keyword_groups):
        if all(kw in text for kw in keywords):
            return 1.0 - 0.1 * i
    return 0


CATEGORY_HEADER_HINTS = [("award", "category"), ("nomination", "category"), ("category",)]
NAME_HEADER_HINTS = [("name", "company"), ("company",), ("nominee", "name"),
                      ("nominee",), ("operator",), ("participant",)]
ZONE_HEADER_HINTS = [("zone",), ("region",)]
RESULT_HEADER_HINTS = [("result",), ("winner",), ("award", "status"), ("status",),
                        ("position",), ("rank",), ("outcome",)]


def detect_columns(ws):
    """Identify the Category / Name / Zone / Result columns by their
    statistical shape in the actual data, NOT by header wording. Header
    text is only used as a small tie-breaking bonus. This keeps the tool
    working even if a future sheet renames every column, as long as the
    underlying data still looks like 'a sparse category column, a dense
    mostly-unique name column, and a dense small-vocabulary result column'
    -- which is true of essentially any award-results spreadsheet.
    """
    header_row = [c.value for c in ws[1]]
    data_rows = list(ws.iter_rows(min_row=2, values_only=True))
    total_rows = len(data_rows)
    ncols = max(len(header_row), max((len(r) for r in data_rows), default=0))

    columns = []
    for i in range(ncols):
        col_vals = [r[i] if i < len(r) else None for r in data_rows]
        columns.append(profile_column(col_vals, total_rows))

    def score(i, structural_fn, hints):
        prof = columns[i]
        if prof["n"] == 0:
            return -1
        base = structural_fn(prof)
        if base < 0:
            return -1
        bonus = header_bonus(header_row[i] if i < len(header_row) else None, hints)
        return base + 0.15 * bonus

    # Category: sparse fill (only the first row of each block has it),
    # fairly long text, high uniqueness among the values it does have.
    def category_score(p):
        if p["fill_ratio"] >= 0.9:
            return -1
        sparsity_score = 1 - p["fill_ratio"]
        length_score = min(p["avg_len"] / 40, 1)
        return sparsity_score * 0.6 + length_score * 0.4 + p["uniq_ratio"] * 0.3

    # Name: dense fill, highly unique, alphabetic-rich, not email/phone-like.
    def name_score(p):
        if p["fill_ratio"] < 0.5 or p["pct_at_sign"] > 0.05 or p["pct_numeric"] > 0.3:
            return -1
        return p["fill_ratio"] * 0.3 + p["uniq_ratio"] * 0.5 + p["pct_alpha_rich"] * 0.2

    # Result/status: dense fill, very LOW distinct-value count relative to
    # row count (small closed vocabulary like Winner/Runner-up), short text.
    def result_score(p):
        if p["fill_ratio"] < 0.5 or p["distinct"] > 8 or p["distinct"] < 2:
            return -1
        closed_vocab_score = 1 - min(p["distinct"] / 8, 1)
        shortness_score = max(0, 1 - p["avg_len"] / 30)
        return p["fill_ratio"] * 0.3 + closed_vocab_score * 0.5 + shortness_score * 0.2

    # Zone/region: same closed-vocabulary idea as result, scored separately
    # so it doesn't get confused with the result column when both exist.
    def zone_score(p):
        if p["fill_ratio"] < 0.5 or p["distinct"] > 10 or p["distinct"] < 2:
            return -1
        closed_vocab_score = 1 - min(p["distinct"] / 10, 1)
        return p["fill_ratio"] * 0.3 + closed_vocab_score * 0.7

    cat_scores = [score(i, category_score, CATEGORY_HEADER_HINTS) for i in range(ncols)]
    name_scores = [score(i, name_score, NAME_HEADER_HINTS) for i in range(ncols)]
    result_scores = [score(i, result_score, RESULT_HEADER_HINTS) for i in range(ncols)]
    zone_scores = [score(i, zone_score, ZONE_HEADER_HINTS) for i in range(ncols)]

    def best(scores, exclude=()):
        candidates = [(s, i) for i, s in enumerate(scores) if i not in exclude]
        if not candidates:
            return None
        s, i = max(candidates)
        return i if s > 0 else None

    name_col = best(name_scores)
    cat_col = best(cat_scores, exclude={name_col} if name_col is not None else set())
    used = {c for c in (name_col, cat_col) if c is not None}
    result_col = best(result_scores, exclude=used)
    used = used | ({result_col} if result_col is not None else set())
    zone_col = best(zone_scores, exclude=used)

    missing = [n for n, v in [("Category", cat_col), ("Company/Nominee Name", name_col),
                               ("Result/Winner status", result_col)] if v is None]
    if missing:
        raise ValueError(
            f"Could not automatically identify the {', '.join(missing)} column(s) from the "
            f"data itself. Header row found: {header_row}. This usually means a needed "
            f"column is missing or its data is too inconsistent to recognize. The sheet "
            f"needs: a column with award category text (filled once per category block), "
            f"a column with company/nominee names (filled on every row), and a column "
            f"with a small set of repeating result labels (e.g. Winner / Runner-up)."
        )

    detected = {
        "category": header_row[cat_col] if cat_col < len(header_row) else f"column {cat_col + 1}",
        "name": header_row[name_col] if name_col < len(header_row) else f"column {name_col + 1}",
        "zone": (header_row[zone_col] if zone_col is not None and zone_col < len(header_row) else None),
        "result": header_row[result_col] if result_col < len(header_row) else f"column {result_col + 1}",
    }
    print(f"Detected columns -> Category: {detected['category']!r}, "
          f"Name: {detected['name']!r}, Zone: {detected['zone']!r}, "
          f"Result: {detected['result']!r}")

    return cat_col, name_col, zone_col, result_col


def determine_winner_index(results):
    """Given the list of result-column values for one category/zone group (in
    original row order), decide which entry is the winner. Two strategies,
    tried in order:

    1. TEXT SIGNAL: if exactly one of the values contains the word "winner"
       (case-insensitive), use that one. This is the common case and is
       cheap/safe to check for literally -- it doesn't depend on the
       *column* being named anything in particular, only on a result VALUE
       containing that word, which is a near-universal convention.
    2. POSITIONAL FALLBACK: if the text signal doesn't cleanly identify
       exactly one winner (e.g. the labels are something else entirely,
       like "1"/"2"/"3", or a Yes/No flag, or just blank), fall back to
       "the first row of the block is the winner" -- which matches how
       these sheets are conventionally built (best result listed first)
       and requires no assumption about wording at all.
    """
    text_matches = [i for i, r in enumerate(results) if r and "winner" in str(r).lower()]
    if len(text_matches) == 1:
        return text_matches[0]
    return 0  # positional fallback: first row in the block


def parse_excel(path):
    wb = openpyxl.load_workbook(path)
    ws = wb.active
    cat_col, name_col, zone_col, result_col = detect_columns(ws)

    rows = list(ws.iter_rows(min_row=2, values_only=True))
    groups = []
    current_category = None
    current_entries = []

    def flush():
        if current_entries:
            results = [e[2] for e in current_entries]
            winner_idx = determine_winner_index(results)
            entries_with_flag = [
                (name, zone, result, i == winner_idx)
                for i, (name, zone, result) in enumerate(current_entries)
            ]
            groups.append({"category": current_category, "entries": entries_with_flag})

    for r in rows:
        cat = r[cat_col] if cat_col < len(r) else None
        name = r[name_col] if name_col < len(r) else None
        zone = r[zone_col] if (zone_col is not None and zone_col < len(r)) else None
        result = r[result_col] if result_col < len(r) else None

        is_blank = all(v is None for v in r)
        if is_blank:
            flush()
            current_entries = []
            continue
        if cat:
            flush()
            current_category = cat
            current_entries = []
        if name:
            current_entries.append((str(name).strip(), zone, result))
    flush()
    return groups


def split_category_title(raw_category):
    """Strip leading numbering ('1.', '2. ', etc.) and split into
    (title, subtitle-or-None) where subtitle is a trailing parenthetical."""
    s = re.sub(r"^\d+\.\s*", "", str(raw_category)).strip()
    m = re.match(r"^(.*?)\s*(\([^)]*\))\s*$", s)
    if m:
        return m.group(1).strip(), m.group(2).strip()
    return s, None


def zone_counts(groups):
    """How many zone-groups exist per category title (to decide whether to
    append a ' - Zone' suffix to the displayed title)."""
    counts = {}
    for g in groups:
        title, _ = split_category_title(g["category"])
        counts[title] = counts.get(title, 0) + 1
    return counts


def clean_company_name(name):
    """Display names exactly as entered in the Excel. Deliberately does NOT
    auto-fix casing (e.g. ALL-CAPS rows) because heuristics can't reliably
    tell a real acronym (e.g. 'ATI', 'ASG') from a normal word typed in caps
    -- and silently mangling a company's own name on the LED screen is worse
    than leaving the organizers' original capitalization untouched. Only
    trims stray whitespace.
    """
    return re.sub(r"\s+", " ", name).strip()


# --------------------------------------------------------------------------
# PPTX template handling
# --------------------------------------------------------------------------

# Literal placeholder tokens. Matching is STRICT and literal: a text box's
# entire content (after stripping the surrounding << >> wrapper and
# whitespace) must equal the token exactly, case-insensitively. No fuzzy
# or partial matching -- "<<winner placeholder>>" does NOT match "WINNER".
TOKEN_NOMINEES = "NOMINEES"            # marks the Nominee stencil + nominee-list box
TOKEN_WINNER = "WINNER"                # marks the Winner stencil + winner-name box
TOKEN_AWARD_CATEGORY = "AWARD CATEGORY"  # optional, on either stencil
TOKEN_ZONE = "ZONE"                    # optional, on either stencil
TOKEN_NOMINEES_WORD = "NOMINEES-WORD"  # optional, anywhere -> literal "NOMINEES"
TOKEN_WINNER_WORD = "WINNER-WORD"      # optional, anywhere -> literal "WINNER"


def shape_text_contains_token(shape, token):
    """True if the shape's text IS the literal token (after stripping
    the surrounding << >> wrapper and whitespace) -- not merely containing
    the token as one word inside a longer sentence. Strict literal match,
    case-insensitive only (e.g. "<<WINNER>>" matches, but
    "<<winner placeholder>>" or "<<Winner Name>>" do NOT).
    """
    if not shape.has_text_frame:
        return False
    stripped = re.sub(r"^[<\[\{\s]+|[>\]\}\s]+$", "", shape.text_frame.text.strip())
    return stripped.upper() == token.upper()


def find_shape_by_token(slide, token):
    for shape in slide.shapes:
        if shape_text_contains_token(shape, token):
            return shape
    return None


def classify_template_slides(template_path):
    """Scan every slide in the template and classify it by which literal
    placeholder tokens it contains:
      - a slide with <<NOMINEES>> is the Nominee stencil
      - a slide with <<WINNER>> is the Winner stencil
      - a slide with neither is left alone (copied through unchanged)
    Every other token (<<AWARD CATEGORY>>, <<ZONE>>, <<nominees-word>>,
    <<winner-word>>) is optional and does not affect classification.
    A template may use ONE of these stencils, or both, in any order.

    Returns a dict: {"nominee": idx_or_None, "winner": idx_or_None}
    """
    prs = Presentation(template_path)
    nominee_idx = None
    winner_idx = None

    for i, slide in enumerate(prs.slides):
        has_nominees = find_shape_by_token(slide, TOKEN_NOMINEES) is not None
        has_winner = find_shape_by_token(slide, TOKEN_WINNER) is not None
        if has_nominees and nominee_idx is None:
            nominee_idx = i
        if has_winner and winner_idx is None:
            winner_idx = i

    return {"nominee": nominee_idx, "winner": winner_idx}


def duplicate_slide(prs, source_slide):
    """Clone a slide (including relationships, e.g. embedded images/EMF
    backgrounds) and append it to the presentation."""
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)

    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for blip in new_slide.shapes._spTree.findall(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
    ):
        old_rid = blip.get(f"{{{r_ns}}}embed")
        if old_rid is None:
            continue
        image_part = source_slide.part.rels[old_rid].target_part
        new_rid = new_slide.part.relate_to(
            image_part, f"{r_ns}/image"
        )
        blip.set(f"{{{r_ns}}}embed", new_rid)

    return new_slide


def find_stencil_shapes(slide, role):
    """Find the optional (award_category_shape, zone_shape) on a stencil
    slide, plus the role-defining names box (nominee list or winner name).
    `role` is 'nominee' or 'winner'. All of these except the names box
    itself may be None -- every placeholder besides the one that defines
    the slide's role is optional."""
    names_shape = find_shape_by_token(slide, TOKEN_NOMINEES if role == "nominee" else TOKEN_WINNER)
    award_category_shape = find_shape_by_token(slide, TOKEN_AWARD_CATEGORY)
    zone_shape = find_shape_by_token(slide, TOKEN_ZONE)
    return award_category_shape, names_shape, zone_shape


def find_shape_by_remembered_name(slide, name):
    if not name:
        return None
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_text_lines(shape, lines):
    """Replace text frame content with one paragraph per line, copying
    formatting (font/size/color/bold) from the existing first paragraph."""
    tf = shape.text_frame
    txBody = tf._txBody
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    paras = txBody.findall(f"{ns_a}p")
    first_p_template = copy.deepcopy(paras[0])
    for p in paras:
        txBody.remove(p)

    for line in lines:
        new_p = copy.deepcopy(first_p_template)
        r_elems = new_p.findall(f"{ns_a}r")
        for extra_r in r_elems[1:]:
            new_p.remove(extra_r)
        t_elem = new_p.find(f"{ns_a}r/{ns_a}t")
        if t_elem is None:
            r_elem = new_p.find(f"{ns_a}r")
            t_elem = r_elem.makeelement(f"{ns_a}t", {})
            r_elem.append(t_elem)
        t_elem.text = line
        txBody.append(new_p)


def set_single_line(shape, text):
    """Set a shape's text frame to a single line of text (used for the
    optional ZONE box and the literal word-label boxes), keeping the
    first paragraph/run's formatting."""
    set_text_lines(shape, [text])


def populate_stencil_copy(prs, stencil_slide, role, title_lines, names, zone_text):
    """Duplicate `stencil_slide`, fill in its placeholder boxes for one
    award (one category/zone, one role: 'nominee' or 'winner'), and
    return the new slide. Every box keeps its own original position, size,
    and text formatting -- only the text content is replaced. Every
    placeholder except the role-defining names box (<<NOMINEES>> or
    <<WINNER>>) is optional and simply skipped if not present."""
    award_category_shape_orig, names_shape_orig, zone_shape_orig = find_stencil_shapes(
        stencil_slide, role
    )
    if names_shape_orig is None:
        token = TOKEN_NOMINEES if role == "nominee" else TOKEN_WINNER
        raise ValueError(f"Could not find <<{token}>> on the {role} stencil slide.")

    # Optional literal word-labels -- work on ANY stencil, independent of
    # that slide's role: <<nominees-word>> always becomes "NOMINEES",
    # <<winner-word>> always becomes "WINNER", wherever they appear.
    nominees_word_shape_orig = find_shape_by_token(stencil_slide, TOKEN_NOMINEES_WORD)
    winner_word_shape_orig = find_shape_by_token(stencil_slide, TOKEN_WINNER_WORD)

    award_category_name = award_category_shape_orig.name if award_category_shape_orig is not None else None
    names_name = names_shape_orig.name
    zone_name = zone_shape_orig.name if zone_shape_orig is not None else None
    nominees_word_name = nominees_word_shape_orig.name if nominees_word_shape_orig is not None else None
    winner_word_name = winner_word_shape_orig.name if winner_word_shape_orig is not None else None

    new_slide = duplicate_slide(prs, stencil_slide)

    award_category_shape = (
        find_shape_by_remembered_name(new_slide, award_category_name) if award_category_name else None
    )
    names_shape = find_shape_by_remembered_name(new_slide, names_name)
    zone_shape = find_shape_by_remembered_name(new_slide, zone_name) if zone_name else None
    nominees_word_shape = (
        find_shape_by_remembered_name(new_slide, nominees_word_name) if nominees_word_name else None
    )
    winner_word_shape = (
        find_shape_by_remembered_name(new_slide, winner_word_name) if winner_word_name else None
    )

    # Names box: the full nominee list, or the single winner's name (left
    # blank if no winner could be identified for this award).
    set_text_lines(names_shape, names if names else [""])

    if award_category_shape is not None:
        set_text_lines(award_category_shape, title_lines)

    if zone_shape is not None:
        if zone_text:
            set_single_line(zone_shape, zone_text)
        else:
            # No zone for this award -- remove the box entirely rather than
            # leaving the literal "<<ZONE>>" token visible on the actual
            # slide (python-pptx has no shape-level visibility toggle, so
            # removing it from the slide's shape tree is the correct way
            # to hide it).
            zone_shape._element.getparent().remove(zone_shape._element)

    if nominees_word_shape is not None:
        set_single_line(nominees_word_shape, "NOMINEES")

    if winner_word_shape is not None:
        set_single_line(winner_word_shape, "WINNER")

    return new_slide


def build_combined_deck(template_path, groups, output_path, seed=RANDOM_SEED):
    """Build ONE deck containing, for every category/zone award: a Nominee
    slide immediately followed by its Winner slide -- using whichever
    slide(s) in the template are tagged with the <<NOMINEES>> / <<WINNER>>
    placeholders. Either one alone is fine (e.g. a Winner-only template
    with no nominee reveal). Any other slide in the template (no
    recognized token) is left exactly as-is and copied through once,
    unchanged, kept at the END of the deck (after every generated award
    slide), since there's no natural "per-award" position for a generic
    divider slide.
    """
    rng = random.Random(seed)
    counts = zone_counts(groups)

    classification = classify_template_slides(template_path)
    nominee_idx = classification["nominee"]
    winner_idx = classification["winner"]

    if nominee_idx is None and winner_idx is None:
        raise ValueError(
            f"Could not find a Nominee stencil (a slide containing <<{TOKEN_NOMINEES}>>) "
            f"or a Winner stencil (a slide containing <<{TOKEN_WINNER}>>) anywhere in the "
            f"template. At least one slide needs a text box whose content is exactly "
            f"'<<{TOKEN_NOMINEES}>>' and/or '<<{TOKEN_WINNER}>>' (literal match -- e.g. "
            f"'<<winner placeholder>>' will NOT match)."
        )

    prs = Presentation(template_path)
    original_slide_count = len(prs.slides)

    nominee_stencil = prs.slides[nominee_idx] if nominee_idx is not None else None
    winner_stencil = prs.slides[winner_idx] if winner_idx is not None else None

    # IMPORTANT: only ever ADD slides during this loop; delete the original
    # stencil/passthrough slides at the very end (see delete_slide's
    # docstring for why interleaving deletes with adds corrupts the file).
    for g in groups:
        title, subtitle = split_category_title(g["category"])
        zone = g["entries"][0][1] if g["entries"] else None
        zone_text = str(zone) if zone else None
        if counts.get(title, 0) > 1 and zone:
            title_with_zone = f"{title} - {zone}"
        else:
            title_with_zone = title
        title_lines = [title_with_zone, subtitle] if subtitle else [title_with_zone]

        if nominee_stencil is not None:
            nominee_names = [clean_company_name(n) for n, z, res, is_win in g["entries"]]
            rng.shuffle(nominee_names)
            populate_stencil_copy(
                prs, nominee_stencil, "nominee", title_lines, nominee_names, zone_text,
            )

        if winner_stencil is not None:
            winners = [clean_company_name(n) for n, z, res, is_win in g["entries"] if is_win]
            # Always generate the Winner slide for this award; if no winner
            # could be identified, the winner box is simply left blank
            # rather than skipping the slide entirely.
            populate_stencil_copy(
                prs, winner_stencil, "winner", title_lines, winners[:1], zone_text,
            )

    # Done adding. Now safe to delete the original stencil slides (see
    # delete_slide's docstring for why this must happen after all adds).
    # Any OTHER original slide (no recognized token -- e.g. a title or
    # thank-you slide) is intentionally kept, but explicitly moved to the
    # end of the deck, after every generated award slide: simply leaving
    # its <p:sldId> entry where it was would put it BEFORE the newly
    # appended slides (since new slides are always appended after
    # whatever already existed), which is not what "left at the end"
    # should mean for a closing/divider slide.
    sld_id_list = prs.slides._sldIdLst
    original_id_elems = list(sld_id_list)[:original_slide_count]
    stencil_elems = [original_id_elems[i] for i in {nominee_idx, winner_idx} - {None}]
    passthrough_elems = [e for e in original_id_elems if e not in stencil_elems]

    for elem in passthrough_elems:
        sld_id_list.remove(elem)
        sld_id_list.append(elem)

    for elem in stencil_elems:
        rId = elem.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        sld_id_list.remove(elem)
        prs.part.drop_rel(rId)

    prs.save(output_path)
    print(f"Saved {output_path} ({len(prs.slides._sldIdLst)} slides)")


def main():
    if len(sys.argv) != 4:
        print("Usage: generate_award_slides.py <excel_file> <template_pptx> <output_file.pptx>")
        sys.exit(1)
    excel_path, template_path, output_path = sys.argv[1], sys.argv[2], sys.argv[3]

    for label, path in [("Excel file", excel_path), ("Template file", template_path)]:
        if not os.path.isfile(path):
            print(f"ERROR: {label} not found at: {path}")
            sys.exit(1)

    out_dir = os.path.dirname(output_path) or "."
    os.makedirs(out_dir, exist_ok=True)

    try:
        groups = parse_excel(excel_path)
        print(f"Parsed {len(groups)} category/zone groups from {excel_path}")
        if not groups:
            print("ERROR: No category/nominee rows were found. Check the Excel file's layout.")
            sys.exit(1)

        build_combined_deck(template_path, groups, output_path)
        print("Done.")
    except ValueError as e:
        print(f"ERROR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
